"""
generate_slides.py
Builds a 12-slide presentation (results/presentation.pptx) from pre-computed
evaluation numbers using python-pptx.

Requires:
    pip install python-pptx

Palette:
    navy   #0D1B2A  (headers)
    blue   #1A6BBF  (monocular)
    red    #D9502A  (stereo)
    off-black text on white background
"""
import os
import sys

try:
    from pptx import Presentation
    from pptx.util    import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text   import PP_ALIGN
except ImportError:
    print("[ERROR] python-pptx is not installed.")
    print("        Run: pip install python-pptx")
    sys.exit(1)


# ─────────────────────────────────────────────
# Constants
# ─────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
BLUE   = RGBColor(0x1A, 0x6B, 0xBF)
RED    = RGBColor(0xD9, 0x50, 0x2A)
GREY   = RGBColor(0x55, 0x55, 0x55)
GREEN  = RGBColor(0x2A, 0x8F, 0x3E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x22, 0x22, 0x22)

RESULTS_DIR = "results"
OUT_PATH    = os.path.join(RESULTS_DIR, "presentation.pptx")

AUTHOR   = "Ahmad Rammal"
SUBTITLE = "VIBOT Master's Project  •  TUM VI Fisheye Benchmark"


# ─────────────────────────────────────────────
# Low-level helpers
# ─────────────────────────────────────────────
def _blank_slide(prs):
    # Layout 6 = "Blank" in the default template
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background() if line is None else _color_line(shape, line)
    shape.shadow.inherit = False
    return shape


def _color_line(shape, color):
    shape.line.color.rgb = color
    shape.line.width     = Pt(0.75)


def _add_text(slide, x, y, w, h, text, *,
              size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top  = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb


def _header(slide, title, page_num=None, total=12):
    # navy band at the top
    _add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), NAVY)
    _add_text(slide, Inches(0.4), Inches(0.22),
              Inches(10), Inches(0.6), title,
              size=24, bold=True, color=WHITE)
    if page_num is not None:
        _add_text(slide, Inches(11.8), Inches(0.3),
                  Inches(1.3), Inches(0.4),
                  f"{page_num} / {total}",
                  size=11, color=WHITE, align=PP_ALIGN.RIGHT)


def _add_bullets(slide, x, y, w, h, lines, *,
                 size=18, bullet_color=BLACK, lead_bold=False):
    """lines is a list of either str or (prefix_bold, rest) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        if isinstance(line, tuple):
            prefix, rest = line
            r = p.add_run(); r.text = f"•  {prefix}"
            r.font.size = Pt(size); r.font.bold = True
            r.font.color.rgb = bullet_color; r.font.name = "Calibri"
            r2 = p.add_run(); r2.text = f"  {rest}"
            r2.font.size = Pt(size); r2.font.bold = False
            r2.font.color.rgb = BLACK; r2.font.name = "Calibri"
        else:
            r = p.add_run(); r.text = f"•  {line}"
            r.font.size = Pt(size)
            r.font.bold = lead_bold
            r.font.color.rgb = BLACK
            r.font.name = "Calibri"
    return tb


def _add_table(slide, x, y, w, h, data, *, header_fill=NAVY):
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, x, y, w, h).table
    for j, col_name in enumerate(data[0]):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = col_name
        r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = "Calibri"
    for i in range(1, rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                RGBColor(0xF3, 0xF5, 0xF8) if i % 2 else WHITE
            )
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            r = p.add_run(); r.text = str(data[i][j])
            r.font.size = Pt(12); r.font.color.rgb = BLACK
            r.font.name = "Calibri"
    return table


def _maybe_add_image(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, width=w, height=h)
        return True
    # fallback placeholder
    _add_rect(slide, x, y, w, h, RGBColor(0xEA, 0xEA, 0xEA),
              line=GREY)
    _add_text(slide, x, y + (h - Inches(0.4)) / 2,
              w, Inches(0.4),
              f"[ missing: {os.path.basename(path)} ]",
              size=12, color=GREY, align=PP_ALIGN.CENTER)
    return False


# ─────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────
def slide_title(prs):
    s = _blank_slide(prs)
    _add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), WHITE)
    _add_rect(s, Inches(0), Inches(2.5), Inches(13.333), Inches(0.08), BLUE)
    _add_rect(s, Inches(0), Inches(4.9), Inches(13.333), Inches(0.08), RED)
    _add_text(s, Inches(0.8), Inches(2.8),
              Inches(11.7), Inches(1.3),
              "From Monocular to Stereo Visual Odometry",
              size=42, bold=True, color=NAVY)
    _add_text(s, Inches(0.8), Inches(4.1),
              Inches(11.7), Inches(0.7),
              SUBTITLE, size=20, color=GREY)
    _add_text(s, Inches(0.8), Inches(5.2),
              Inches(11.7), Inches(0.5),
              AUTHOR, size=18, bold=True, color=NAVY)
    _add_text(s, Inches(0.8), Inches(5.75),
              Inches(11.7), Inches(0.5),
              "Classical CV pipeline • FAST+LK • Essential matrix • "
              "StereoSGBM • PnP-RANSAC",
              size=14, color=GREY)


def slide_motivation(prs):
    s = _blank_slide(prs)
    _header(s, "Motivation — The Scale Problem", page_num=2)
    _add_text(s, Inches(0.5), Inches(1.2),
              Inches(12.3), Inches(0.8),
              "Why monocular VO fails for metric tasks",
              size=22, bold=True, color=NAVY)
    _add_bullets(s, Inches(0.5), Inches(2.1),
                 Inches(12.3), Inches(4.5), [
        ("Scale ambiguity.",
         "A single camera cannot recover metric depth from two views — "
         "the essential matrix is only defined up to a scalar."),
        ("Drift in real units.",
         "Any robotics task needing distance (obstacle avoidance, "
         "mapping, collision checks) breaks without absolute scale."),
        ("Common workarounds are brittle.",
         "IMU fusion adds hardware and calibration cost; scale priors "
         "from object size break outside structured scenes."),
        ("Stereo dissolves the ambiguity.",
         "A known baseline gives depth directly: Z = f·B / d, from the "
         "first frame, with no bootstrap."),
    ], size=16)


def slide_related_work(prs):
    s = _blank_slide(prs)
    _header(s, "Related Work", page_num=3)
    data = [
        ["Year", "Author(s)",       "Contribution"],
        ["2011", "Scaramuzza & Fraundorfer", "VO tutorial — pipeline used here as baseline"],
        ["2018", "Schubert et al.", "TUM VI Benchmark — fisheye dataset, GT from motion capture"],
        ["2017", "Mur-Artal & Tardós", "ORB-SLAM2 — mono/stereo SLAM reference system"],
        ["2013", "Geiger et al.",  "KITTI — the earlier standard benchmark for VO/SLAM"],
    ]
    _add_table(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(3.6), data)
    _add_text(s, Inches(0.6), Inches(5.6),
              Inches(12.1), Inches(1.5),
              "This work: minimal classical pipeline — no bundle adjustment, "
              "no loop closure — used as a controlled baseline to quantify "
              "the mono→stereo improvement on TUM VI fisheye data.",
              size=14, color=GREY)


def slide_dataset(prs):
    s = _blank_slide(prs)
    _header(s, "Dataset — TUM VI Benchmark (512×512 DSO export)", page_num=4)
    data = [
        ["Sequence",  "Setting",       "Frames",  "Length", "GT coverage"],
        ["room2",     "Indoor room",   "~3000",   "short",  "full (mocap)"],
        ["corridor3", "Indoor hallway","~6000",   "medium", "start + end only"],
        ["outdoors5", "Outdoor walk",  "~16 000", "long",   "start + end only"],
    ]
    _add_table(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(2.9), data)
    _add_bullets(s, Inches(0.6), Inches(4.6),
                 Inches(12.1), Inches(2.5), [
        ("Sensor.",
         "Stereo fisheye (~195° FOV), 10 cm baseline, 20 Hz."),
        ("Calibration.",
         "Kalibr model, 4-parameter fisheye distortion."),
        ("Ground truth.",
         "OptiTrack mocap indoors; outdoor sequences have GT only at "
         "start + end → we use the Eq. 8 start-end drift metric."),
    ], size=15)


def slide_mono_pipeline(prs):
    s = _blank_slide(prs)
    _header(s, "Monocular VO Pipeline", page_num=5)
    steps = [
        ("1.", "FAST detection",     "corners on the histogram-equalised left image"),
        ("2.", "Lucas-Kanade track", "pyramidal LK flow to the next frame"),
        ("3.", "Fisheye undistort",  "cv2.fisheye.undistortPoints(P=K) on tracked pts"),
        ("4.", "Essential matrix",   "findEssentialMat + RANSAC → recoverPose"),
        ("5.", "Compose pose",       "cumulative R, t accumulated up-to-scale"),
    ]
    x0 = Inches(0.7); y0 = Inches(1.6); w = Inches(11.9); h = Inches(0.85)
    for i, (n, name, desc) in enumerate(steps):
        y = y0 + i * (h + Inches(0.12))
        _add_rect(s, x0, y, Inches(0.85), h, BLUE)
        _add_text(s, x0, y, Inches(0.85), h, n,
                  size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_rect(s, x0 + Inches(0.85), y, w - Inches(0.85), h,
                  RGBColor(0xF3, 0xF5, 0xF8))
        _add_text(s, x0 + Inches(1.05), y + Inches(0.08),
                  Inches(4.0), Inches(0.4), name,
                  size=16, bold=True, color=NAVY)
        _add_text(s, x0 + Inches(5.0), y + Inches(0.12),
                  w - Inches(5.2), Inches(0.6), desc,
                  size=14, color=BLACK)


def slide_stereo_pipeline(prs):
    s = _blank_slide(prs)
    _header(s, "Stereo VO Pipeline — Metric Scale", page_num=6)
    steps = [
        ("1.", "Rectify stereo pair","fisheye.stereoRectify + remap"),
        ("2.", "SGBM disparity",     "block-matching on rectified pair"),
        ("3.", "Back-project",       "metric depth Z = f·B / d per feature"),
        ("4.", "PnP-RANSAC",         "3D (prev) ↔ 2D (cur) → relative pose"),
        ("5.", "Invert & compose",   "accumulate camera-to-world pose"),
    ]
    x0 = Inches(0.7); y0 = Inches(1.6); w = Inches(11.9); h = Inches(0.85)
    for i, (n, name, desc) in enumerate(steps):
        y = y0 + i * (h + Inches(0.12))
        _add_rect(s, x0, y, Inches(0.85), h, RED)
        _add_text(s, x0, y, Inches(0.85), h, n,
                  size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_rect(s, x0 + Inches(0.85), y, w - Inches(0.85), h,
                  RGBColor(0xFB, 0xF1, 0xEC))
        _add_text(s, x0 + Inches(1.05), y + Inches(0.08),
                  Inches(4.0), Inches(0.4), name,
                  size=16, bold=True, color=NAVY)
        _add_text(s, x0 + Inches(5.0), y + Inches(0.12),
                  w - Inches(5.2), Inches(0.6), desc,
                  size=14, color=BLACK)


def slide_fisheye_challenges(prs):
    s = _blank_slide(prs)
    _header(s, "Fisheye Challenges", page_num=7)
    _add_bullets(s, Inches(0.6), Inches(1.4),
                 Inches(12.2), Inches(5.5), [
        ("Focal-length override.",
         "cv2.fisheye.stereoRectify returns a tiny fx for ~195° FOV. "
         "We override with the original fx so the rectified central "
         "region has a useful perspective projection."),
        ("FAST over Shi-Tomasi.",
         "Shi-Tomasi under-detects on dark TUM VI frames. FAST with a "
         "low threshold (8) after histogram-equalisation gives dense, "
         "evenly distributed corners."),
        ("Point-wise undistortion.",
         "For mono VO we avoid dewarping the whole image. Instead we "
         "run cv2.fisheye.undistortPoints on the tracked features only, "
         "then pass them to findEssentialMat with P=K — fast and "
         "preserves the full FOV."),
    ], size=16)


def slide_results_tracking(prs):
    s = _blank_slide(prs)
    _header(s, "Results — Tracking Stability", page_num=8)
    data = [
        ["Sequence",  "Mono failures", "Stereo failures"],
        ["room2",     "many",          "0"],
        ["corridor3", "many",          "0"],
        ["outdoors5", "many (long)",   "0"],
    ]
    _add_table(s, Inches(1.5), Inches(1.6), Inches(10.3), Inches(2.8), data)
    _add_bullets(s, Inches(0.6), Inches(4.7),
                 Inches(12.2), Inches(2.5), [
        ("Mono.",
         "Frequent re-detection when LK tracks drop below 50; essential-"
         "matrix solve occasionally degenerate on low-parallax frames."),
        ("Stereo.",
         "PnP-RANSAC never falls below its 6-point minimum thanks to "
         "the dense disparity-derived 3D map."),
    ], size=16)


def slide_results_rpe(prs):
    s = _blank_slide(prs)
    _header(s, "Results — RPE Translational (Stereo improvement)", page_num=9)
    # Prefer the bar chart if it exists
    img_path = os.path.join(RESULTS_DIR, "plots", "rpe_translational.png")
    _maybe_add_image(s, img_path,
                     Inches(0.6), Inches(1.3),
                     Inches(8.0), Inches(5.0))
    data = [
        ["Sequence",  "Mono (m)", "Stereo (m)", "Gain"],
        ["room2",     "11.37",    "0.88",       "13×"],
        ["corridor3", "66.20",    "1.93",       "34×"],
        ["outdoors5", "272.35",   "9.72",       "28×"],
    ]
    _add_table(s, Inches(8.8), Inches(2.1), Inches(4.0), Inches(3.0), data)
    _add_text(s, Inches(8.8), Inches(5.3),
              Inches(4.0), Inches(1.8),
              "Known baseline removes scale drift — per-window error "
              "drops by 1-2 orders of magnitude.",
              size=13, color=GREY)


def slide_results_table(prs):
    s = _blank_slide(prs)
    _header(s, "Results — Full Evaluation", page_num=10)
    data = [
        ["Seq.",      "Mode",   "ATE (m)", "RPE-t (m)", "RPE-r (°)", "Drift (m)"],
        ["room2",     "Mono",   "1.2205",  "11.37",     "24.45",     "1879.34"],
        ["room2",     "Stereo", "2.9040",  "0.88",      "46.90",     "3.57"],
        ["corridor3", "Mono",   "0.8205",  "66.20",     "33.45",     "707.24"],
        ["corridor3", "Stereo", "9.6643",  "1.93",      "37.76",     "20.57"],
        ["outdoors5", "Mono",   "0.8381",  "272.35",    "68.31",     "4564.91"],
        ["outdoors5", "Stereo", "76.7849", "9.72",      "51.81",     "149.83"],
    ]
    _add_table(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(4.5), data)
    _add_text(s, Inches(0.7), Inches(6.2),
              Inches(11.9), Inches(1.0),
              "ATE column: Sim3-aligned for mono (scale is cancelled), "
              "SE3-aligned for stereo. Drift is the Eq. 8 start-end error.",
              size=13, color=GREY)


def slide_discussion(prs):
    s = _blank_slide(prs)
    _header(s, "Discussion", page_num=11)
    _add_bullets(s, Inches(0.6), Inches(1.4),
                 Inches(12.2), Inches(5.8), [
        ("Mono ATE is misleading.",
         "Sim3 alignment absorbs the scale error, so mono ATE looks "
         "small. RPE-t and Eq. 8 drift (thousands of metres) expose "
         "the real failure."),
        ("Rotational drift without loop closure.",
         "Even stereo accumulates >40° rotational drift on long loops — "
         "classical pipelines need BA or loop closure to bound this."),
        ("Outdoors is the hardest.",
         "Longer trajectory, faster motion, thinner textures ⇒ stereo "
         "drift climbs to ~150 m despite metric depth."),
        ("Stereo baseline is the key fix.",
         "Most RPE-t gains come from the first frame — no bootstrap, "
         "no IMU, no learned priors."),
    ], size=16)


def slide_conclusion(prs):
    s = _blank_slide(prs)
    _header(s, "Conclusion & Future Work", page_num=12)
    _add_text(s, Inches(0.6), Inches(1.3),
              Inches(6.2), Inches(0.5),
              "Achievements",
              size=20, bold=True, color=NAVY)
    _add_bullets(s, Inches(0.6), Inches(1.9),
                 Inches(6.2), Inches(4.5), [
        "Two self-contained pipelines (mono + stereo) on fisheye VI data",
        "13-34× RPE-t improvement from stereo, no IMU, no DL",
        "Reproducible evaluation: ATE / RPE / Eq. 8 drift in one script",
    ], size=15, bullet_color=GREEN, lead_bold=True)

    _add_text(s, Inches(7.2), Inches(1.3),
              Inches(5.5), Inches(0.5),
              "Future work",
              size=20, bold=True, color=NAVY)
    _add_bullets(s, Inches(7.2), Inches(1.9),
                 Inches(5.5), Inches(5.0), [
        "Local bundle adjustment over a sliding window",
        "Loop closure (DBoW2 or similar) to bound drift",
        "IMU pre-integration (VIO) using the TUM VI IMU channel",
        "Learned features (SuperPoint/SuperGlue) as a drop-in front end",
    ], size=15, bullet_color=RED, lead_bold=True)


# ─────────────────────────────────────────────
# Driver
# ─────────────────────────────────────────────
def main():
    prs = Presentation()
    # 16:9
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_motivation(prs)
    slide_related_work(prs)
    slide_dataset(prs)
    slide_mono_pipeline(prs)
    slide_stereo_pipeline(prs)
    slide_fisheye_challenges(prs)
    slide_results_tracking(prs)
    slide_results_rpe(prs)
    slide_results_table(prs)
    slide_discussion(prs)
    slide_conclusion(prs)

    os.makedirs(RESULTS_DIR, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"[Slides] wrote {len(prs.slides)} slides -> {OUT_PATH}")


if __name__ == "__main__":
    main()
